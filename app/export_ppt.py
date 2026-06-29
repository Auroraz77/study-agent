from __future__ import annotations

import re
from collections.abc import Iterable
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Microsoft YaHei"
MONO_FONT = "Consolas"

INK = RGBColor(31, 41, 55)
MUTED = RGBColor(102, 112, 133)
BLUE = RGBColor(37, 99, 235)
NAVY = RGBColor(30, 64, 175)
GREEN = RGBColor(15, 118, 110)
AMBER = RGBColor(180, 83, 9)
RED = RGBColor(185, 28, 28)
SURFACE = RGBColor(248, 250, 252)
LINE = RGBColor(217, 222, 231)
LIGHT_BLUE = RGBColor(239, 246, 255)
LIGHT_GREEN = RGBColor(236, 253, 245)
LIGHT_AMBER = RGBColor(255, 247, 237)
LIGHT_RED = RGBColor(254, 242, 242)
DARK = RGBColor(16, 24, 40)


def build_learning_pptx(payload: dict[str, Any], student_name: str | None = None) -> BytesIO:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    course = _clean_text(payload.get("course")) or "Learning Plan"
    resources = payload.get("resources") if isinstance(payload.get("resources"), list) else []
    profile = payload.get("profile") if isinstance(payload.get("profile"), dict) else {}
    path = payload.get("learning_path") if isinstance(payload.get("learning_path"), dict) else {}

    _add_title_slide(
        prs,
        course,
        "Personalized learning resource deck",
        [f"Student: {student_name}" if student_name else "", f"Resources: {len(resources)}"],
    )
    _add_summary_slide(prs, payload)
    if profile:
        _add_profile_slide(prs, profile)
    if path:
        _add_path_slide(prs, path)

    for index, resource in enumerate(resources, start=1):
        if not isinstance(resource, dict):
            continue
        _add_resource_intro_slide(prs, resource, index, len(resources))
        if resource.get("type") == "quiz" and isinstance(resource.get("quiz"), list):
            _add_quiz_slides(prs, resource)
        else:
            _add_markdown_resource_slides(prs, resource)

    if len(prs.slides) == 1:
        _add_content_slide(prs, "No generated resources", ["Generate a learning plan first, then export again."])

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out


def suggested_filename(course: str | None) -> str:
    base = _clean_text(course) or "learning-plan"
    base = re.sub(r'[\\/:*?"<>|]+', "_", base).strip(" ._") or "learning-plan"
    return f"{base[:80]}.pptx"


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_round_rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _add_rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def _add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    size: int = 16,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN | None = None,
    font: str = FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _clean_text(text)
    if align is not None:
        p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def _add_title_slide(prs: Presentation, title: str, subtitle: str, meta: Iterable[str] = ()) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, LIGHT_BLUE)
    _add_accent_bar(slide, BLUE)

    _add_round_rect(slide, Inches(0.85), Inches(0.85), Inches(11.9), Inches(5.75), RGBColor(255, 255, 255), LINE)
    _add_text(slide, _clean_text(title), Inches(1.25), Inches(1.7), Inches(10.9), Inches(1.4), 36, NAVY, True)
    _add_rect(slide, Inches(1.3), Inches(3.22), Inches(3.2), Inches(0.05), BLUE)
    _add_text(slide, subtitle, Inches(1.28), Inches(3.45), Inches(8.8), Inches(0.55), 18, MUTED)

    clean_meta = [item for item in meta if item]
    if clean_meta:
        _add_text(slide, "  |  ".join(clean_meta), Inches(1.28), Inches(5.8), Inches(8.8), Inches(0.5), 13, MUTED)


def _add_summary_slide(prs: Presentation, payload: dict[str, Any]) -> None:
    final_answer = _clean_text(payload.get("final_answer"))
    resources = payload.get("resources") if isinstance(payload.get("resources"), list) else []
    lines = _wrap_lines(final_answer, 5) if final_answer else []
    if not lines:
        lines = [f"Generated {len(resources)} personalized resources."]
    _add_content_slide(prs, "Learning plan summary", lines[:5], accent=BLUE, style="cards")


def _add_profile_slide(prs: Presentation, profile: dict[str, Any]) -> None:
    cards = []
    for key, value in profile.items():
        rendered = " / ".join(str(item) for item in value) if isinstance(value, list) else str(value)
        if rendered.strip():
            cards.append((str(key), rendered))
    _add_key_value_cards_slide(prs, "Student profile", cards[:6] if cards else [("Profile", "No profile data.")])


def _add_path_slide(prs: Presentation, path: dict[str, Any]) -> None:
    title = _clean_text(path.get("title")) or "Learning path"
    stages = path.get("stages") if isinstance(path.get("stages"), list) else []
    timeline = []
    for index, stage in enumerate(stages, start=1):
        if not isinstance(stage, dict):
            continue
        name = _clean_text(stage.get("name")) or f"Stage {index}"
        goal = _clean_text(stage.get("goal"))
        resources = stage.get("resources")
        resource_text = " / ".join(str(item) for item in resources) if isinstance(resources, list) else ""
        timeline.append({"index": index, "name": name, "goal": goal, "resources": resource_text})
    if not timeline:
        _add_content_slide(prs, title, ["No path stages."], accent=GREEN, style="cards")
        return
    for idx, chunk in enumerate(_chunk_texts(timeline, 4), start=1):
        _add_timeline_slide(prs, title if idx == 1 else f"{title} ({idx})", chunk)


def _add_resource_intro_slide(prs: Presentation, resource: dict[str, Any], index: int, total: int) -> None:
    title = _clean_text(resource.get("title")) or f"Resource {index}"
    agent = _clean_text(resource.get("agent"))
    modality = _clean_text(resource.get("modality"))
    bullets = [f"Resource {index} of {total}", f"Type: {_clean_text(resource.get('type')) or 'resource'}"]
    if agent:
        bullets.append(f"Agent: {agent}")
    if modality:
        bullets.append(f"Modality: {modality}")
    _add_content_slide(prs, title, bullets, accent=GREEN, style="hero")


def _add_markdown_resource_slides(prs: Presentation, resource: dict[str, Any]) -> None:
    title = _clean_text(resource.get("title")) or "Resource"
    sections = _parse_markdown_sections(str(resource.get("content") or ""))
    if not sections:
        sections = [(title, [{"type": "text", "text": _clean_text(resource.get("content")) or "No content."}])]

    for section_title, blocks in sections:
        clean_title = section_title if section_title and section_title != title else title
        pending: list[str] = []
        for block in blocks:
            if block.get("type") == "code":
                if pending:
                    for chunk in _chunk_texts(pending, 5):
                        _add_content_slide(prs, clean_title, chunk, style="cards")
                    pending = []
                _add_code_slide(prs, clean_title, block.get("text", ""))
            elif block.get("type") == "table":
                if pending:
                    for chunk in _chunk_texts(pending, 5):
                        _add_content_slide(prs, clean_title, chunk, style="cards")
                    pending = []
                _add_table_as_bullets(prs, clean_title, block.get("rows", []))
            else:
                pending.extend(_wrap_lines(block.get("text", ""), 2))
        if pending:
            for chunk in _chunk_texts(pending, 5):
                _add_content_slide(prs, clean_title, chunk, style="cards")


def _add_quiz_slides(prs: Presentation, resource: dict[str, Any]) -> None:
    title = _clean_text(resource.get("title")) or "Quiz"
    quiz = [item for item in resource.get("quiz", []) if isinstance(item, dict)]
    for index, item in enumerate(quiz, start=1):
        _add_quiz_question_slide(prs, title, item, index)


def _add_content_slide(
    prs: Presentation,
    title: str,
    lines: list[str],
    accent: RGBColor = BLUE,
    style: str = "cards",
) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, RGBColor(255, 255, 255))
    _add_header(slide, title, accent)

    clean_lines = [_truncate(_clean_text(line), 96) for line in (lines or [""]) if _clean_text(line)]
    if not clean_lines:
        clean_lines = ["No content."]

    if style == "hero":
        _add_round_rect(slide, Inches(0.9), Inches(1.65), Inches(11.5), Inches(4.8), SURFACE, LINE)
        for idx, line in enumerate(clean_lines[:4]):
            top = Inches(2.05 + idx * 0.85)
            _add_badge(slide, Inches(1.25), top, str(idx + 1), accent)
            _add_text(slide, line, Inches(1.85), top - Inches(0.03), Inches(9.8), Inches(0.45), 18, INK, idx == 0)
        _add_footer(slide)
        return

    card_count = min(len(clean_lines), 5)
    card_h = Inches(0.82 if card_count >= 5 else 0.95)
    gap = Inches(0.18)
    start_y = Inches(1.5)
    for idx, line in enumerate(clean_lines[:5]):
        top = start_y + idx * (card_h + gap)
        fill = [LIGHT_BLUE, LIGHT_GREEN, LIGHT_AMBER, SURFACE, LIGHT_RED][idx % 5]
        _add_round_rect(slide, Inches(0.9), top, Inches(11.45), card_h, fill, LINE)
        _add_badge(slide, Inches(1.18), top + Inches(0.22), str(idx + 1), accent)
        _add_text(slide, line, Inches(1.75), top + Inches(0.17), Inches(10.1), card_h - Inches(0.18), 16, INK)
    _add_footer(slide)


def _add_code_slide(prs: Presentation, title: str, code: str) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, RGBColor(255, 255, 255))
    _add_header(slide, title, BLUE)

    _add_round_rect(slide, Inches(0.8), Inches(1.35), Inches(8.15), Inches(5.65), DARK, DARK)
    _add_round_rect(slide, Inches(9.25), Inches(1.35), Inches(3.1), Inches(5.65), LIGHT_BLUE, LINE)
    _add_text(slide, "Run & observe", Inches(9.55), Inches(1.75), Inches(2.5), Inches(0.4), 18, NAVY, True)
    _add_text(
        slide,
        "Focus on the input, the model step, and the final output. Use the code as a runnable checkpoint.",
        Inches(9.55),
        Inches(2.35),
        Inches(2.35),
        Inches(1.7),
        14,
        MUTED,
    )
    _add_text(slide, "Challenge", Inches(9.55), Inches(4.4), Inches(2.5), Inches(0.35), 16, GREEN, True)
    _add_text(slide, "Change one parameter and compare the result.", Inches(9.55), Inches(4.9), Inches(2.35), Inches(0.8), 14, INK)

    box = slide.shapes.add_textbox(Inches(1.05), Inches(1.6), Inches(7.65), Inches(5.2))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    lines = str(code or "").strip().splitlines()
    numbered = [f"{idx + 1:02d}  {line}" for idx, line in enumerate(lines[:22])]
    p.text = "\n".join(numbered) or "No code."
    p.font.name = MONO_FONT
    p.font.size = Pt(10.5)
    p.font.color.rgb = RGBColor(248, 250, 252)
    _add_footer(slide)


def _add_table_as_bullets(prs: Presentation, title: str, rows: list[list[str]]) -> None:
    if not rows:
        return
    header = rows[0]
    data = rows[1:6]
    if not data:
        _add_content_slide(prs, title, ["No table rows."])
        return

    slide = _blank_slide(prs)
    _set_bg(slide, RGBColor(255, 255, 255))
    _add_header(slide, title, BLUE)
    cols = min(len(header), max(len(row) for row in data), 4)
    table_shape = slide.shapes.add_table(len(data) + 1, cols, Inches(0.8), Inches(1.55), Inches(11.75), Inches(4.8))
    table = table_shape.table
    for col in range(cols):
        table.cell(0, col).text = _truncate(header[col] if col < len(header) else "", 32)
        _style_table_cell(table.cell(0, col), LIGHT_BLUE, NAVY, True, 12)
    for row_idx, row in enumerate(data, start=1):
        for col in range(cols):
            table.cell(row_idx, col).text = _truncate(row[col] if col < len(row) else "", 70)
            _style_table_cell(table.cell(row_idx, col), RGBColor(255, 255, 255), INK, False, 11)
    _add_footer(slide)


def _add_key_value_cards_slide(prs: Presentation, title: str, cards: list[tuple[str, str]]) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, RGBColor(255, 255, 255))
    _add_header(slide, title, BLUE)
    positions = [
        (0.85, 1.55), (4.65, 1.55), (8.45, 1.55),
        (0.85, 3.85), (4.65, 3.85), (8.45, 3.85),
    ]
    fills = [LIGHT_BLUE, LIGHT_GREEN, LIGHT_AMBER, SURFACE, LIGHT_RED, RGBColor(245, 243, 255)]
    for idx, (key, value) in enumerate(cards[:6]):
        x, y = positions[idx]
        _add_round_rect(slide, Inches(x), Inches(y), Inches(3.35), Inches(1.75), fills[idx], LINE)
        _add_text(slide, _truncate(key, 24), Inches(x + 0.28), Inches(y + 0.22), Inches(2.8), Inches(0.35), 14, NAVY, True)
        _add_text(slide, _truncate(value, 82), Inches(x + 0.28), Inches(y + 0.72), Inches(2.75), Inches(0.82), 13, INK)
    _add_footer(slide)


def _add_timeline_slide(prs: Presentation, title: str, stages: list[dict[str, Any]]) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, RGBColor(255, 255, 255))
    _add_header(slide, title, GREEN)
    y = Inches(1.65)
    for idx, stage in enumerate(stages):
        top = y + idx * Inches(1.22)
        color = [GREEN, BLUE, AMBER, NAVY][idx % 4]
        _add_badge(slide, Inches(0.95), top + Inches(0.22), str(stage.get("index", idx + 1)), color)
        if idx < len(stages) - 1:
            _add_rect(slide, Inches(1.17), top + Inches(0.72), Inches(0.04), Inches(0.75), LINE)
        _add_round_rect(slide, Inches(1.65), top, Inches(10.55), Inches(0.98), SURFACE, LINE)
        _add_text(slide, _truncate(stage.get("name", ""), 36), Inches(1.95), top + Inches(0.12), Inches(3.1), Inches(0.32), 15, color, True)
        _add_text(slide, _truncate(stage.get("goal", ""), 86), Inches(5.05), top + Inches(0.1), Inches(6.7), Inches(0.34), 13, INK)
        if stage.get("resources"):
            _add_text(slide, _truncate(f"Resources: {stage['resources']}", 96), Inches(5.05), top + Inches(0.5), Inches(6.7), Inches(0.28), 11, MUTED)
    _add_footer(slide)


def _add_quiz_question_slide(prs: Presentation, title: str, item: dict[str, Any], index: int) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, RGBColor(255, 255, 255))
    _add_header(slide, f"{title} - Q{index}", AMBER)

    question = _clean_text(item.get("question")) or f"Question {index}"
    _add_round_rect(slide, Inches(0.85), Inches(1.32), Inches(11.65), Inches(1.15), LIGHT_AMBER, RGBColor(253, 186, 116))
    _add_badge(slide, Inches(1.15), Inches(1.62), "Q", AMBER)
    _add_text(slide, _truncate(question, 105), Inches(1.75), Inches(1.55), Inches(10.05), Inches(0.5), 17, INK, True)

    options = item.get("options") if isinstance(item.get("options"), list) else []
    option_positions = [(0.9, 2.85), (6.75, 2.85), (0.9, 3.85), (6.75, 3.85)]
    for opt_idx, option in enumerate(options[:4]):
        x, y = option_positions[opt_idx]
        _add_round_rect(slide, Inches(x), Inches(y), Inches(5.4), Inches(0.72), SURFACE, LINE)
        _add_badge(slide, Inches(x + 0.2), Inches(y + 0.17), chr(65 + opt_idx), BLUE)
        _add_text(slide, _truncate(_strip_option_label(str(option)), 55), Inches(x + 0.75), Inches(y + 0.16), Inches(4.25), Inches(0.36), 13, INK)

    answer = _clean_text(item.get("answer"))
    explanation = _clean_text(item.get("explanation"))
    _add_round_rect(slide, Inches(0.9), Inches(5.3), Inches(11.55), Inches(1.0), LIGHT_GREEN, RGBColor(167, 243, 208))
    if answer:
        _add_text(slide, f"Answer: {answer}", Inches(1.2), Inches(5.52), Inches(2.1), Inches(0.3), 14, GREEN, True)
    if explanation:
        _add_text(slide, _truncate(explanation, 120), Inches(3.2), Inches(5.48), Inches(8.8), Inches(0.42), 13, INK)
    _add_footer(slide)


def _add_badge(slide, x, y, text: str, color: RGBColor):
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.38), Inches(0.38))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.color.rgb = color
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def _style_table_cell(cell, fill: RGBColor, color: RGBColor, bold: bool, size: int) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def _add_footer(slide) -> None:
    _add_text(slide, "Personalized Learning Deck", Inches(0.85), Inches(7.08), Inches(4.2), Inches(0.22), 9, MUTED)
    _add_rect(slide, Inches(10.75), Inches(7.16), Inches(1.65), Inches(0.025), LINE)


def _add_header(slide, title: str, accent: RGBColor) -> None:
    _add_accent_bar(slide, accent)
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.38), Inches(11.8), Inches(0.65))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = _clean_text(title)[:90]
    p.font.name = FONT
    p.font.bold = True
    p.font.size = Pt(25)
    p.font.color.rgb = INK
    _add_rect(slide, Inches(0.78), Inches(1.12), Inches(11.65), Inches(0.03), accent)


def _add_accent_bar(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _parse_markdown_sections(markdown: str) -> list[tuple[str, list[dict[str, Any]]]]:
    sections: list[tuple[str, list[dict[str, Any]]]] = []
    current_title = "Content"
    blocks: list[dict[str, Any]] = []
    lines = str(markdown or "").replace("\r\n", "\n").split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            i += 1
            continue
        heading = re.match(r"^(#{1,4})\s+(.+)$", stripped)
        if heading:
            if blocks:
                sections.append((current_title, blocks))
                blocks = []
            current_title = _clean_markdown_inline(heading.group(2))
            i += 1
            continue
        if stripped.startswith("```"):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1
            blocks.append({"type": "code", "text": "\n".join(code_lines)})
            continue
        if _is_table_start(lines, i):
            table_lines = [lines[i], lines[i + 1]]
            i += 2
            while i < len(lines) and "|" in lines[i]:
                table_lines.append(lines[i])
                i += 1
            blocks.append({"type": "table", "rows": _parse_table(table_lines)})
            continue
        if re.match(r"^([-*+]|\d+\.)\s+", stripped):
            items = []
            while i < len(lines) and re.match(r"^([-*+]|\d+\.)\s+", lines[i].strip()):
                items.append(_clean_markdown_inline(re.sub(r"^([-*+]|\d+\.)\s+", "", lines[i].strip())))
                i += 1
            blocks.extend({"type": "text", "text": item} for item in items if item)
            continue

        paragraph = [stripped]
        i += 1
        while i < len(lines) and lines[i].strip() and not _is_block_start(lines, i):
            paragraph.append(lines[i].strip())
            i += 1
        blocks.append({"type": "text", "text": _clean_markdown_inline(" ".join(paragraph))})

    if blocks:
        sections.append((current_title, blocks))
    return sections


def _is_block_start(lines: list[str], index: int) -> bool:
    stripped = lines[index].strip()
    return bool(
        re.match(r"^(#{1,4})\s+", stripped)
        or stripped.startswith("```")
        or re.match(r"^([-*+]|\d+\.)\s+", stripped)
        or _is_table_start(lines, index)
    )


def _is_table_start(lines: list[str], index: int) -> bool:
    return (
        index + 1 < len(lines)
        and "|" in lines[index]
        and re.match(r"^\s*\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?\s*$", lines[index + 1])
        is not None
    )


def _parse_table(lines: list[str]) -> list[list[str]]:
    return [_split_table_row(line) for idx, line in enumerate(lines) if idx != 1]


def _split_table_row(line: str) -> list[str]:
    return [_clean_markdown_inline(cell.strip()) for cell in line.strip().strip("|").split("|")]


def _chunk_texts(lines: list[str], max_items: int) -> list[list[str]]:
    chunks = []
    current = []
    for line in lines:
        if len(current) >= max_items:
            chunks.append(current)
            current = []
        current.append(line)
    if current:
        chunks.append(current)
    return chunks or [[]]


def _wrap_lines(text: str, max_lines: int) -> list[str]:
    clean = _clean_text(text)
    if not clean:
        return []
    parts = re.split(r"(?<=[。！？.!?])\s+|\n+", clean)
    lines = [part.strip() for part in parts if part.strip()]
    if len(lines) <= max_lines:
        return lines
    return lines[: max_lines - 1] + [_truncate(" ".join(lines[max_lines - 1 :]), 160)]


def _clean_markdown_inline(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = text.replace("$", "")
    return _clean_text(text)


def _clean_text(value: Any) -> str:
    text = str(value or "")
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _truncate(text: str, limit: int) -> str:
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "..."


def _strip_option_label(option: str) -> str:
    return re.sub(r"^\s*[\(（]?[A-Ga-g][\)）]?\s*(?:[.．、:：]\s*)?", "", option).strip()
