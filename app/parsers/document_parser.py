from __future__ import annotations

from io import BytesIO
from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


class ParseError(Exception):
    pass


def parse_document(filename: str, data: bytes) -> str:
    suffix = Path(filename or "").suffix.lower()

    if suffix in {".txt", ".md", ".csv"}:
        return _parse_text(data)
    if suffix == ".pdf":
        return _parse_pdf(data)
    if suffix == ".docx":
        return _parse_docx(data)
    if suffix == ".pptx":
        return _parse_pptx(data)

    raise ParseError(f"暂不支持解析该文件类型：{suffix or 'unknown'}")


def _parse_text(data: bytes) -> str:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ParseError("文本文件不是 UTF-8 编码") from exc
    return _require_text(text)


def _parse_pdf(data: bytes) -> str:
    try:
        reader = PdfReader(BytesIO(data))
        pages: list[str] = []
        for index, page in enumerate(reader.pages, start=1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages.append(f"第 {index} 页\n{page_text.strip()}")
    except Exception as exc:
        raise ParseError(f"PDF 解析失败：{exc}") from exc

    return _require_text(
        "\n\n".join(pages),
        "PDF 未提取到可用文本，可能是扫描版 PDF，需要 OCR。",
    )


def _parse_docx(data: bytes) -> str:
    try:
        document = Document(BytesIO(data))
        parts: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)

        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
    except Exception as exc:
        raise ParseError(f"Word 文档解析失败：{exc}") from exc

    return _require_text("\n".join(parts), "Word 文档未提取到可用文本。")


def _parse_pptx(data: bytes) -> str:
    try:
        presentation = Presentation(BytesIO(data))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        texts.append(text)
            if texts:
                slides.append(f"第 {index} 页幻灯片\n" + "\n".join(texts))
    except Exception as exc:
        raise ParseError(f"PPT 文档解析失败：{exc}") from exc

    return _require_text("\n\n".join(slides), "PPT 文档未提取到可用文本。")


def _require_text(text: str, message: str = "文件未提取到可用文本。") -> str:
    normalized = text.strip()
    if not normalized:
        raise ParseError(message)
    return normalized
